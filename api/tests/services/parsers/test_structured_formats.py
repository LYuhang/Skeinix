"""Knowledge parser coverage for office, web, JSON and tabular sources."""
from __future__ import annotations

import json
from io import BytesIO

from openpyxl import Workbook
from pptx import Presentation

from vibecanvas_api.services.parsers import detect_parser_type
from vibecanvas_api.services.parsers.csv import CsvParser
from vibecanvas_api.services.parsers.html import HtmlParser
from vibecanvas_api.services.parsers.json import JsonParser
from vibecanvas_api.services.parsers.pptx import PptxParser
from vibecanvas_api.services.parsers.xlsx import XlsxParser


def test_detects_supported_structured_sources() -> None:
    assert detect_parser_type("rows.csv", "text/csv") == "csv"
    assert detect_parser_type("data.json", "application/json") == "json"
    assert detect_parser_type("page.html", "text/html") == "html"
    assert detect_parser_type(
        "deck.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ) == "pptx"
    assert detect_parser_type(
        "table.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ) == "xlsx"
    assert detect_parser_type("legacy.doc", "application/msword") is None


def test_csv_rows_carry_row_metadata() -> None:
    segments = CsvParser().parse(b"name,team\nAda,Platform\nLin,Research\n")
    assert segments[0].metadata == {"row": 2}
    assert "name: Ada" in segments[0].text
    assert "team: Research" in segments[1].text


def test_json_records_carry_json_paths() -> None:
    blob = json.dumps([{"name": "Ada"}, {"name": "Lin"}]).encode()
    segments = JsonParser().parse(blob)
    assert [segment.metadata["json_path"] for segment in segments] == ["$[0]", "$[1]"]
    assert "$[1].name: Lin" in segments[1].text


def test_html_discards_scripts_and_preserves_heading() -> None:
    segments = HtmlParser().parse(
        b"<html><body><h1>Policy</h1><p>Keep this.</p>"
        b"<script>ignore me</script></body></html>"
    )
    assert segments[0].metadata == {"heading": "Policy"}
    assert "Keep this" in segments[0].text
    assert "ignore me" not in segments[0].text


def test_xlsx_preserves_sheet_and_rows() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "People"
    sheet.append(["name", "team"])
    sheet.append(["Ada", "Platform"])
    buffer = BytesIO()
    workbook.save(buffer)
    segments = XlsxParser().parse(buffer.getvalue())
    assert segments[0].metadata == {"sheet": "People", "row_start": 1, "row_end": 2}
    assert "Ada | Platform" in segments[0].text


def test_pptx_preserves_slide_number() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Ship Knowledge retrieval"
    buffer = BytesIO()
    presentation.save(buffer)
    segments = PptxParser().parse(buffer.getvalue())
    assert segments[0].metadata == {"slide": 1}
    assert "Roadmap" in segments[0].text
