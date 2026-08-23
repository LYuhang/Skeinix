from __future__ import annotations

import json
from pathlib import Path
from types import SimpleNamespace

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas

from vibecanvas_api.document_runtime import render as render_module
from vibecanvas_api.document_runtime.render import render_document_feedback
from vibecanvas_api.document_runtime.review import review_document
from vibecanvas_api.document_runtime.server import mcp


def test_document_mcp_registers_structured_review_and_feedback_tools() -> None:
    tools = {tool.name: tool for tool in mcp._tool_manager.list_tools()}

    assert set(tools) == {"review_document", "render_document_feedback"}
    assert tools["review_document"].output_schema is not None
    assert tools["render_document_feedback"].output_schema is not None
    feedback_schema = tools["render_document_feedback"].parameters
    assert feedback_schema["properties"]["dpi"] == {
        "default": 144,
        "maximum": 220,
        "minimum": 96,
        "title": "Dpi",
        "type": "integer",
    }
    assert feedback_schema["properties"]["max_pages"] == {
        "default": 8,
        "maximum": 20,
        "minimum": 1,
        "title": "Max Pages",
        "type": "integer",
    }


def test_review_document_covers_native_office_and_delivery_formats(
    tmp_path: Path,
) -> None:
    docx_path = tmp_path / "brief.docx"
    document = Document()
    document.add_heading("Quarterly brief", level=1)
    document.add_paragraph("A concise decision memo.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Owner"
    table.cell(1, 0).text = "Operations"
    document.save(docx_path)

    pptx_path = tmp_path / "review.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Executive review"
    slide.shapes.add_textbox(0, 0, 2_000_000, 500_000).text = "Decision"
    deck.save(pptx_path)

    xlsx_path = tmp_path / "forecast.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Forecast"
    sheet.append(["Month", "Revenue", "Cost", "Margin"])
    sheet.append(["Jan", 120, 70, "=B2-C2"])
    sheet.freeze_panes = "A2"
    sheet.auto_filter.ref = "A1:D2"
    workbook.save(xlsx_path)

    pdf_path = tmp_path / "invoice.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 760, "Invoice 2026-08")
    pdf.showPage()
    pdf.save()

    docx_result = review_document(str(docx_path))
    assert docx_result["valid"] is True
    assert docx_result["details"]["headings"] == ["Quarterly brief"]
    assert docx_result["details"]["tables"] == 1

    pptx_result = review_document(str(pptx_path))
    assert pptx_result["valid"] is True
    assert pptx_result["details"]["slide_count"] == 1

    xlsx_result = review_document(str(xlsx_path))
    assert xlsx_result["valid"] is True
    assert xlsx_result["details"]["sheets"][0]["formulas"] == 1
    assert xlsx_result["details"]["sheets"][0]["freeze_panes"] == "A2"

    pdf_result = review_document(str(pdf_path))
    assert pdf_result["valid"] is True
    assert pdf_result["details"]["page_count"] == 1
    assert pdf_result["details"]["sampled_page_text_characters"][0] > 0


def test_review_document_covers_lightweight_office_companion_formats(
    tmp_path: Path,
) -> None:
    samples = {
        "handbook.md": "# Handbook\n\n## Procedure\n\n1. Review\n2. Publish\n",
        "dashboard.html": (
            "<!doctype html><html><head><title>Dashboard</title></head>"
            "<body><h1>Operations</h1><h2>Risk</h2></body></html>"
        ),
        "events.csv": "timestamp,status\n2026-08-22,ok\n",
        "notes.txt": "Release notes\n",
        "visual.svg": (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 60">'
            '<rect width="100" height="60" fill="#eef2ff"/></svg>'
        ),
    }
    results: dict[str, dict] = {}
    for name, content in samples.items():
        path = tmp_path / name
        path.write_text(content, encoding="utf-8")
        results[name] = review_document(str(path))

    assert all(result["valid"] for result in results.values())
    assert results["handbook.md"]["details"]["headings"] == [
        "# Handbook",
        "## Procedure",
    ]
    assert results["dashboard.html"]["details"]["title"] == "Dashboard"
    assert results["events.csv"]["details"]["rows"] == 2
    assert results["visual.svg"]["details"]["view_box"] == "0 0 100 60"


def test_render_feedback_reports_only_actual_omission(
    tmp_path: Path,
    monkeypatch,
) -> None:
    pdf_path = tmp_path / "report.pdf"
    pdf_path.write_bytes(b"%PDF-test")
    feedback_root = tmp_path / "feedback"
    monkeypatch.setattr(render_module, "_FEEDBACK_ROOT", feedback_root)
    monkeypatch.setattr(render_module, "_required_command", lambda *_: "pdftoppm")

    def fake_run(arguments, **_kwargs):
        prefix = Path(arguments[-1])
        for page_number in range(1, 4):
            prefix.with_name(f"{prefix.name}-{page_number}.png").write_bytes(
                b"png" + bytes([page_number])
            )
        return SimpleNamespace(returncode=0, stdout="", stderr="")

    monkeypatch.setattr(render_module.subprocess, "run", fake_run)

    result = render_document_feedback(str(pdf_path), max_pages=2)

    assert result["rendered_pages"] == 2
    assert result["truncated"] is True
    assert len(result["feedback_paths"]) == 2
    assert all(Path(path).is_file() for path in result["feedback_paths"])
    assert json.dumps(result).count("page-003") == 0
