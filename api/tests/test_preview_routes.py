from __future__ import annotations

from io import BytesIO
import uuid

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import pytest
from sqlalchemy import text

from vibecanvas_api.services.chat_workspace import chat_workspace_scope_id
from vibecanvas_api.services.file_revision import vfs_content_revision
from vibecanvas_api.services.object_store import get_object_store
from vibecanvas_api.storage.chat_repo import ChatRepo
from vibecanvas_api.storage.db import session_scope
from vibecanvas_api.storage.vfs_store import VfsRepo


def _office_payload(extension: str) -> bytes:
    target = BytesIO()
    if extension == "docx":
        document = Document()
        document.add_heading("Executive brief", level=1)
        document.add_paragraph("Decision-ready content")
        document.save(target)
    elif extension == "pptx":
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Executive review"
        presentation.save(target)
    else:
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Summary"
        sheet.append(["Metric", "Value"])
        sheet.append(["Adoption", 0.72])
        workbook.save(target)
    return target.getvalue()


async def _register(client) -> tuple[dict, dict]:
    response = await client.post(
        "/api/v1/auth/register",
        json={
            "email": f"preview_{uuid.uuid4().hex[:12]}@example.com",
            "username": "Preview User",
            "password": "pw12345678",
        },
    )
    assert response.status_code in (200, 201), response.text
    headers = {"Authorization": f"Bearer {response.json()['session_token']}"}
    me = (await client.get("/api/v1/auth/me", headers=headers)).json()
    return headers, me


async def _chat_fixture(client, app_engine) -> tuple[dict, dict, str, str]:
    headers, me = await _register(client)
    boot = await client.get("/api/v1/chats/bootstrap?surface=chat", headers=headers)
    chat_id = f"c_preview_{uuid.uuid4().hex[:10]}"
    async with session_scope(tenant_id=me["tenant_id"]) as session:
        await ChatRepo(session, me["user_id"]).register_session(
            boot.json()["carrier_scope_id"],
            name="Preview",
            chat_id=chat_id,
            surface="chat",
        )
    return headers, me, chat_id, chat_workspace_scope_id(chat_id)


async def _upload(client, headers, scope_id: str, name: str, data: bytes, mime: str):
    response = await client.post(
        f"/api/v1/vfs/upload?wf_id={scope_id}&folder=data",
        headers=headers,
        files={"file": (name, data, mime)},
    )
    assert response.status_code == 200, response.text
    return response.json()["path"]


async def _write_chat_workspace_file(
    *, tenant_id: str, scope_id: str, path: str, data: bytes, mime: str
) -> None:
    async with session_scope(tenant_id=tenant_id) as session:
        await VfsRepo(session, object_store=get_object_store()).upsert_artifact_bytes(
            wf_id=scope_id,
            tenant=tenant_id,
            path=path,
            data=data,
            content_type=mime,
        )


async def _artifact_events(
    app_engine,
    *,
    tenant_id: str,
    scope_id: str,
    path: str,
) -> list[dict]:
    async with app_engine.begin() as connection:
        await connection.execute(
            text("SELECT set_config('app.tenant_id',:tenant,true)"),
            {"tenant": tenant_id},
        )
        rows = (
            await connection.execute(
                text(
                    "SELECT event_id, event_type, content_revision "
                    "FROM vfs_artifact_events "
                    "WHERE scope_kind = 'artifact' "
                    "AND scope_id = :scope_id AND path = :path "
                    "ORDER BY event_id"
                ),
                {"scope_id": scope_id, "path": path},
            )
        ).mappings().all()
    return [dict(row) for row in rows]


@pytest.mark.asyncio
async def test_native_drawio_resolves_to_official_viewer_descriptor(
    client, app_engine,
):
    headers, _me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    source = b'''<mxGraphModel><root>
      <mxCell id="0"/><mxCell id="1" parent="0"/>
      <mxCell id="start" value="Start" vertex="1" parent="1">
        <mxGeometry x="40" y="40" width="120" height="60" as="geometry"/>
      </mxCell>
      <mxCell id="end" value="End" vertex="1" parent="1">
        <mxGeometry x="260" y="40" width="120" height="60" as="geometry"/>
      </mxCell>
      <mxCell id="flow" edge="1" source="start" target="end" parent="1">
        <mxGeometry relative="1" as="geometry"/>
      </mxCell>
    </root></mxGraphModel>'''
    path = await _upload(
        client,
        headers,
        scope_id,
        "native.drawio",
        source,
        "application/vnd.jgraph.mxfile",
    )
    response = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={"fileRef": {
            "schemaVersion": 1,
            "scope": "chat",
            "chatId": chat_id,
            "path": path,
        }},
    )
    assert response.status_code == 200, response.text
    descriptor = response.json()
    assert descriptor["renderer"] == "drawio"
    assert descriptor["contentType"] == "application/vnd.jgraph.mxfile"
    assert descriptor["loadPolicy"] == "range"
    assert descriptor["diagram"]["status"] == "valid"
    assert descriptor["diagram"]["summary"] == {
        "cells": 5,
        "vertices": 2,
        "edges": 1,
        "pages": 1,
    }
    content = await client.get(descriptor["content"]["url"])
    assert content.status_code == 200
    assert content.content == source


@pytest.mark.asyncio
async def test_preview_text_save_preserves_bom_newlines_and_revision(
    client, app_engine,
):
    headers, me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    path = await _upload(
        client,
        headers,
        scope_id,
        "notes.md",
        b"\xef\xbb\xbf# Title\r\n\r\nBefore\r\n",
        "text/markdown",
    )
    file_ref = {
        "schemaVersion": 1,
        "scope": "chat",
        "chatId": chat_id,
        "path": path,
    }
    resolved = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={"fileRef": file_ref},
    )
    assert resolved.status_code == 200, resolved.text
    descriptor = resolved.json()
    assert descriptor["schemaVersion"] == 1
    assert descriptor["renderer"] == "markdown"
    assert descriptor["capabilities"] == {
        "preview": True,
        "edit": True,
        "download": True,
    }
    assert descriptor["text"] == {
        "encoding": "utf-8",
        "bom": True,
        "newline": "CRLF",
        "mixedNewlines": False,
    }
    assert descriptor["content"]["inlineText"].endswith("Before\r\n")

    initial_events = await _artifact_events(
        app_engine,
        tenant_id=me["tenant_id"],
        scope_id=scope_id,
        path=path,
    )
    assert len(initial_events) == 1
    assert initial_events[0]["event_type"] == "upsert"
    assert (
        vfs_content_revision(initial_events[0]["content_revision"])
        == descriptor["revision"]
    )

    # A VFS read touches last_access for LRU bookkeeping. It must neither
    # change optimistic content identity nor produce a live Preview event.
    read = await client.get(
        "/api/v1/vfs/content",
        headers=headers,
        params={"wf_id": scope_id, "path": path},
    )
    assert read.status_code == 200, read.text
    after_read = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={"fileRef": file_ref},
    )
    assert after_read.status_code == 200, after_read.text
    assert after_read.json()["revision"] == descriptor["revision"]
    assert await _artifact_events(
        app_engine,
        tenant_id=me["tenant_id"],
        scope_id=scope_id,
        path=path,
    ) == initial_events

    saved = await client.put(
        "/api/v1/previews/file",
        headers=headers,
        json={
            "fileRef": file_ref,
            "expectedRevision": descriptor["revision"],
            "contentType": "text/markdown",
            "content": "# Title\n\nAfter\n",
        },
    )
    assert saved.status_code == 200, saved.text
    assert saved.json()["revision"] != descriptor["revision"]

    stale = await client.put(
        "/api/v1/previews/file",
        headers=headers,
        json={
            "fileRef": file_ref,
            "expectedRevision": descriptor["revision"],
            "contentType": "text/markdown",
            "content": "stale",
        },
    )
    assert stale.status_code == 409
    assert stale.json()["detail"] == "preview_revision_conflict"

    signed = await client.post(
        "/api/v1/vfs/sign",
        headers=headers,
        json={"wf_id": scope_id, "path": path},
    )
    raw = await client.get(signed.json()["url"])
    assert raw.content == b"\xef\xbb\xbf# Title\r\n\r\nAfter\r\n"
    assert raw.headers["etag"].startswith('"sha256:')


@pytest.mark.asyncio
async def test_structured_text_tables_are_preview_only(client, app_engine):
    headers, _me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    path = await _upload(
        client,
        headers,
        scope_id,
        "records.csv",
        b"name,value\nalpha,1\n",
        "text/csv",
    )
    file_ref = {
        "schemaVersion": 1,
        "scope": "chat",
        "chatId": chat_id,
        "path": path,
    }
    resolved = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={"fileRef": file_ref},
    )
    assert resolved.status_code == 200, resolved.text
    descriptor = resolved.json()
    assert descriptor["renderer"] == "spreadsheet"
    assert descriptor["capabilities"] == {
        "preview": True,
        "edit": False,
        "download": True,
    }

    write = await client.put(
        "/api/v1/previews/file",
        headers=headers,
        json={
            "fileRef": file_ref,
            "expectedRevision": descriptor["revision"],
            "contentType": "text/csv",
            "content": "name,value\nbeta,2\n",
        },
    )
    assert write.status_code == 403
    assert write.json()["detail"] == "preview_file_read_only"


@pytest.mark.asyncio
@pytest.mark.parametrize(
    ("extension", "mime", "renderer"),
    [
        (
            "docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "docx",
        ),
        (
            "pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx",
        ),
    ],
)
async def test_native_office_preview_uses_authorized_pdf_rendition(
    client,
    app_engine,
    monkeypatch,
    extension,
    mime,
    renderer,
):
    from vibecanvas_api.routes import previews as preview_routes

    headers, _me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    payload = _office_payload(extension)
    path = await _upload(
        client,
        headers,
        scope_id,
        f"brief.{extension}",
        payload,
        mime,
    )
    file_ref = {
        "schemaVersion": 1,
        "scope": "chat",
        "chatId": chat_id,
        "path": path,
    }
    resolved = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={"fileRef": file_ref},
    )
    assert resolved.status_code == 200, resolved.text
    descriptor = resolved.json()
    assert descriptor["renderer"] == renderer
    assert descriptor["rendition"] == {
        "format": "pdf",
        "contentType": "application/pdf",
        "url": descriptor["rendition"]["url"],
        "sourceRevision": descriptor["revision"],
    }
    assert descriptor["rendition"]["url"].startswith(
        "/api/v1/previews/office-rendition?"
    )

    original = await client.get(descriptor["content"]["url"])
    assert original.status_code == 200
    assert original.content == payload

    rendered_pdf = b"%PDF-1.7\nfaithful-office-preview"
    calls = []

    def fake_render(data: bytes, suffix: str) -> bytes:
        calls.append((data, suffix))
        return rendered_pdf

    monkeypatch.setattr(preview_routes, "render_office_preview_pdf", fake_render)
    unauthenticated = await client.get(descriptor["rendition"]["url"])
    assert unauthenticated.status_code == 401

    other_headers, _other = await _register(client)
    forbidden = await client.get(
        descriptor["rendition"]["url"],
        headers=other_headers,
    )
    assert forbidden.status_code == 404

    rendition = await client.get(
        descriptor["rendition"]["url"],
        headers=headers,
    )
    assert rendition.status_code == 200, rendition.text
    assert rendition.headers["content-type"] == "application/pdf"
    assert rendition.headers["x-content-type-options"] == "nosniff"
    assert rendition.content == rendered_pdf
    assert calls == [(payload, f".{extension}")]


@pytest.mark.asyncio
async def test_xlsx_preview_uses_native_workbook_source_without_pdf_rendition(
    client,
    app_engine,
):
    headers, _me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    payload = _office_payload("xlsx")
    path = await _upload(
        client,
        headers,
        scope_id,
        "budget.xlsx",
        payload,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    resolved = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={
            "fileRef": {
                "schemaVersion": 1,
                "scope": "chat",
                "chatId": chat_id,
                "path": path,
            }
        },
    )
    assert resolved.status_code == 200, resolved.text
    descriptor = resolved.json()
    assert descriptor["renderer"] == "spreadsheet"
    assert descriptor["loadPolicy"] == "inline"
    assert descriptor["content"]["url"]
    assert descriptor["rendition"] is None

    original = await client.get(descriptor["content"]["url"])
    assert original.status_code == 200
    assert original.content == payload

@pytest.mark.asyncio
@pytest.mark.parametrize("path", ["/memory/state.md", "/logs/runtime.log"])
async def test_agent_owned_chat_files_are_read_only_previewable(
    client, app_engine, path,
):
    headers, me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    await _write_chat_workspace_file(
        tenant_id=me["tenant_id"],
        scope_id=scope_id,
        path=path,
        data=b"agent-owned context\n",
        mime="text/markdown" if path.endswith(".md") else "text/plain",
    )
    file_ref = {
        "schemaVersion": 1,
        "scope": "chat",
        "chatId": chat_id,
        "path": path,
    }

    resolved = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={"fileRef": file_ref},
    )
    assert resolved.status_code == 200, resolved.text
    descriptor = resolved.json()
    assert descriptor["content"]["inlineText"] == "agent-owned context\n"
    assert descriptor["capabilities"] == {
        "preview": True,
        "edit": False,
        "download": True,
    }

    write = await client.put(
        "/api/v1/previews/file",
        headers=headers,
        json={
            "fileRef": file_ref,
            "expectedRevision": descriptor["revision"],
            "contentType": descriptor["contentType"],
            "content": "user overwrite",
        },
    )
    assert write.status_code == 403
    assert write.json()["detail"] == "preview_file_read_only"

@pytest.mark.asyncio
async def test_preview_detects_pdf_streams_range_and_rejects_plain_archives(
    client, app_engine,
):
    headers, _me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    pdf_path = await _upload(
        client,
        headers,
        scope_id,
        "report.bin",
        b"%PDF-1.7\n" + b"x" * 128,
        "application/octet-stream",
    )
    pdf_ref = {
        "schemaVersion": 1,
        "scope": "chat",
        "chatId": chat_id,
        "path": pdf_path,
    }
    pdf = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={"fileRef": pdf_ref},
    )
    assert pdf.status_code == 200, pdf.text
    assert pdf.json()["renderer"] == "pdf"
    assert pdf.json()["loadPolicy"] == "range"
    ranged = await client.get(
        pdf.json()["content"]["url"],
        headers={"Range": "bytes=0-7"},
    )
    assert ranged.status_code == 206, ranged.text
    assert ranged.content == b"%PDF-1.7"
    assert ranged.headers["content-range"].startswith("bytes 0-7/")

    archive_path = await _upload(
        client,
        headers,
        scope_id,
        "bundle.zip",
        b"PK\x03\x04not-a-preview",
        "application/zip",
    )
    archive = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={
            "fileRef": {
                "schemaVersion": 1,
                "scope": "chat",
                "chatId": chat_id,
                "path": archive_path,
            }
        },
    )
    assert archive.status_code == 200, archive.text
    assert archive.json()["renderer"] == "unsupported"
    assert archive.json()["loadPolicy"] == "unsupported"
    assert archive.json()["error"] == {
        "code": "archive_preview_not_supported",
        "params": {},
    }


@pytest.mark.asyncio
async def test_preview_file_ref_is_owner_scoped(client, app_engine):
    owner_headers, _owner, chat_id, scope_id = await _chat_fixture(client, app_engine)
    path = await _upload(
        client,
        owner_headers,
        scope_id,
        "private.txt",
        b"secret",
        "text/plain",
    )
    other_headers, _other = await _register(client)
    response = await client.post(
        "/api/v1/previews/resolve",
        headers=other_headers,
        json={
            "fileRef": {
                "schemaVersion": 1,
                "scope": "chat",
                "chatId": chat_id,
                "path": path,
            }
        },
    )
    assert response.status_code == 404
    assert response.json()["detail"] == "preview_file_not_found"


@pytest.mark.asyncio
async def test_preview_html_resource_session_maps_workspace_files(
    client, app_engine,
):
    headers, _me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    html_path = await _upload(
        client,
        headers,
        scope_id,
        "index.html",
        b'<img src="/data/pixel.png">',
        "text/html",
    )
    image = b"\x89PNG\r\n\x1a\npreview-resource"
    await _upload(
        client,
        headers,
        scope_id,
        "pixel.png",
        image,
        "image/png",
    )
    await _upload(
        client,
        headers,
        scope_id,
        "private.txt",
        b"not referenced by the HTML",
        "text/plain",
    )
    file_ref = {
        "schemaVersion": 1,
        "scope": "chat",
        "chatId": chat_id,
        "path": html_path,
    }
    response = await client.post(
        "/api/v1/previews/resource-session",
        headers=headers,
        json={"fileRef": file_ref},
    )
    assert response.status_code == 200, response.text
    session = response.json()
    root = next(
        mount["rootUrl"]
        for mount in session["resourceMounts"]
        if mount["pathPrefix"] == "/"
    )
    rendered_resource = await client.get(f"{root}data/pixel.png")
    assert rendered_resource.status_code == 200, rendered_resource.text
    assert rendered_resource.content == image
    unrelated = await client.get(f"{root}data/private.txt")
    assert unrelated.status_code == 403
    wrong_audience = await client.get(
        f"{root}data/pixel.png".replace(
            "/resources/file-preview/",
            "/resources/interactive-artifact/",
        )
    )
    assert wrong_audience.status_code == 403


@pytest.mark.asyncio
async def test_preview_markdown_resource_session_maps_relative_images(
    client, app_engine,
):
    headers, me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    markdown_path = "/data/handbooks/operations-handbook.md"
    image_path = "/data/handbooks/handbook-architecture.svg"
    image = b'<svg xmlns="http://www.w3.org/2000/svg"/>'
    await _write_chat_workspace_file(
        tenant_id=me["tenant_id"],
        scope_id=scope_id,
        path=markdown_path,
        data=b"![Architecture](handbook-architecture.svg)",
        mime="text/markdown",
    )
    await _write_chat_workspace_file(
        tenant_id=me["tenant_id"],
        scope_id=scope_id,
        path=image_path,
        data=image,
        mime="image/svg+xml",
    )
    await _write_chat_workspace_file(
        tenant_id=me["tenant_id"],
        scope_id=scope_id,
        path="/data/handbooks/private.txt",
        data=b"not referenced by the Markdown",
        mime="text/plain",
    )
    response = await client.post(
        "/api/v1/previews/resource-session",
        headers=headers,
        json={
            "fileRef": {
                "schemaVersion": 1,
                "scope": "chat",
                "chatId": chat_id,
                "path": markdown_path,
            }
        },
    )
    assert response.status_code == 200, response.text
    base_url = response.json()["baseUrl"]
    rendered_resource = await client.get(f"{base_url}handbook-architecture.svg")
    assert rendered_resource.status_code == 200, rendered_resource.text
    assert rendered_resource.content == image
    unrelated = await client.get(f"{base_url}private.txt")
    assert unrelated.status_code == 403


@pytest.mark.asyncio
async def test_legacy_office_is_download_only_without_conversion(
    client, app_engine,
):
    headers, _me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    path = await _upload(
        client,
        headers,
        scope_id,
        "legacy.rtf",
        b"{\\rtf1\\ansi Preview conversion}",
        "application/rtf",
    )
    file_ref = {
        "schemaVersion": 1,
        "scope": "chat",
        "chatId": chat_id,
        "path": path,
    }
    response = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={"fileRef": file_ref},
    )
    assert response.status_code == 200, response.text
    descriptor = response.json()
    assert descriptor["renderer"] == "unsupported"
    assert descriptor["loadPolicy"] == "unsupported"
    assert descriptor["capabilities"]["download"] is True
    assert descriptor["error"] == {
        "code": "unsupported_file_type",
        "params": {"extension": ".rtf"},
    }


@pytest.mark.asyncio
async def test_large_table_returns_structured_preview_error(client, app_engine):
    headers, _me, chat_id, scope_id = await _chat_fixture(client, app_engine)
    data = b"name,value\n" + b"a,1\n" * (3 * 1024 * 1024)
    assert len(data) > 10 * 1024 * 1024
    path = await _upload(
        client,
        headers,
        scope_id,
        "large.csv",
        data,
        "text/csv",
    )
    response = await client.post(
        "/api/v1/previews/resolve",
        headers=headers,
        json={
            "fileRef": {
                "schemaVersion": 1,
                "scope": "chat",
                "chatId": chat_id,
                "path": path,
            }
        },
    )
    assert response.status_code == 200, response.text
    descriptor = response.json()
    assert descriptor["renderer"] == "unsupported"
    assert descriptor["loadPolicy"] == "unsupported"
    assert descriptor["error"]["code"] == "file_too_large"
    assert descriptor["error"]["params"]["actualBytes"] == len(data)
    assert descriptor["error"]["params"]["limitBytes"] == 10 * 1024 * 1024
