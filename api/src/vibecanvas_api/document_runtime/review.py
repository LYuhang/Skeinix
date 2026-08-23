"""Deterministic structural checks for files produced by ``/document``.

The checks deliberately cover facts that code can establish reliably. Visual
quality remains a pixel-level Agent review performed on images returned by the
companion renderer.
"""

from __future__ import annotations

import csv
import hashlib
import io
from pathlib import Path
from typing import Any

from vibecanvas_api.services.parsers.archive import validate_office_archive


_MAX_SOURCE_BYTES = 64 * 1024 * 1024
_MAX_MATERIALIZED_SPREADSHEET_CELLS = 1_000_000
_WORKSPACE_ROOTS = ("/data", "/mount", "/memory", "/tmp")


class DocumentReviewError(ValueError):
    """The requested document cannot be reviewed safely."""


def _resolve_workspace_file(value: str) -> Path:
    raw = str(value or "").strip()
    if not raw or "\x00" in raw:
        raise DocumentReviewError("path is required")
    candidate = Path(raw)
    if not candidate.is_absolute():
        candidate = Path.cwd() / candidate
    try:
        resolved = candidate.resolve(strict=True)
    except OSError as exc:
        raise DocumentReviewError(f"document does not exist: {raw}") from exc
    if not resolved.is_file():
        raise DocumentReviewError(f"document is not a file: {raw}")
    allowed_roots = [Path(root).resolve() for root in _WORKSPACE_ROOTS]
    if not any(
        resolved == root or root in resolved.parents
        for root in allowed_roots
    ):
        raise DocumentReviewError(
            "document must be inside the current sandbox workspace"
        )
    size = resolved.stat().st_size
    if size > _MAX_SOURCE_BYTES:
        raise DocumentReviewError("document exceeds the 64 MiB review limit")
    return resolved


def _presentation(blob: bytes) -> tuple[dict[str, Any], list[str], list[str]]:
    from pptx import Presentation

    validate_office_archive(blob, required_prefix="ppt/")
    deck = Presentation(io.BytesIO(blob))
    errors: list[str] = []
    warnings: list[str] = []
    slides: list[dict[str, Any]] = []
    explicit_font_sizes: list[float] = []
    for slide_number, slide in enumerate(deck.slides, start=1):
        outside: list[str] = []
        text_chars = 0
        for index, shape in enumerate(slide.shapes, start=1):
            left = int(getattr(shape, "left", 0) or 0)
            top = int(getattr(shape, "top", 0) or 0)
            width = int(getattr(shape, "width", 0) or 0)
            height = int(getattr(shape, "height", 0) or 0)
            if (
                left < 0
                or top < 0
                or left + width > int(deck.slide_width)
                or top + height > int(deck.slide_height)
            ):
                outside.append(str(getattr(shape, "name", "") or f"shape-{index}"))
            if getattr(shape, "has_text_frame", False):
                text = str(getattr(shape, "text", "") or "")
                text_chars += len(text.strip())
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            explicit_font_sizes.append(float(run.font.size.pt))
        if outside:
            errors.append(
                f"slide {slide_number} has objects outside the page: "
                + ", ".join(outside[:8])
            )
        if not slide.shapes:
            warnings.append(f"slide {slide_number} is empty")
        slides.append({
            "number": slide_number,
            "shapes": len(slide.shapes),
            "text_characters": text_chars,
            "outside_page_objects": outside,
        })
    if not slides:
        errors.append("presentation contains no slides")
    very_small = [size for size in explicit_font_sizes if size < 10]
    if very_small:
        warnings.append(
            f"presentation contains {len(very_small)} explicitly sized text runs below 10 pt"
        )
    return ({
        "slides": slides,
        "slide_count": len(slides),
        "slide_size_emu": {
            "width": int(deck.slide_width),
            "height": int(deck.slide_height),
        },
        "explicit_font_size_pt": {
            "minimum": min(explicit_font_sizes) if explicit_font_sizes else None,
            "maximum": max(explicit_font_sizes) if explicit_font_sizes else None,
        },
    }, errors, warnings)


def _word(blob: bytes) -> tuple[dict[str, Any], list[str], list[str]]:
    from docx import Document

    validate_office_archive(blob, required_prefix="word/")
    document = Document(io.BytesIO(blob))
    errors: list[str] = []
    warnings: list[str] = []
    nonempty = [p for p in document.paragraphs if p.text.strip()]
    headings = [
        p.text.strip()
        for p in nonempty
        if str(getattr(p.style, "name", "") or "").lower().startswith("heading")
    ]
    if not nonempty and not document.tables and not document.inline_shapes:
        errors.append("Word document has no readable content")
    if len(nonempty) >= 12 and not headings:
        warnings.append("long Word document has no heading styles")
    sections = []
    for index, section in enumerate(document.sections, start=1):
        sections.append({
            "number": index,
            "page_width_emu": int(section.page_width),
            "page_height_emu": int(section.page_height),
            "margins_emu": {
                "top": int(section.top_margin),
                "right": int(section.right_margin),
                "bottom": int(section.bottom_margin),
                "left": int(section.left_margin),
            },
        })
    return ({
        "paragraphs": len(document.paragraphs),
        "nonempty_paragraphs": len(nonempty),
        "headings": headings[:50],
        "tables": len(document.tables),
        "inline_images": len(document.inline_shapes),
        "sections": sections,
    }, errors, warnings)


def _spreadsheet(blob: bytes) -> tuple[dict[str, Any], list[str], list[str]]:
    import openpyxl

    validate_office_archive(blob, required_prefix="xl/")
    workbook = openpyxl.load_workbook(io.BytesIO(blob), data_only=False)
    errors: list[str] = []
    warnings: list[str] = []
    sheets = []
    for sheet in workbook.worksheets:
        # ``iter_rows`` walks the full declared rectangle. A workbook with one
        # accidental cell at XFD1048576 would therefore turn a quick review
        # into billions of visits. The materialized-cell map contains exactly
        # the cells that were present in the OOXML worksheet.
        materialized_cells = list(sheet._cells.values())
        if len(materialized_cells) > _MAX_MATERIALIZED_SPREADSHEET_CELLS:
            errors.append(
                f"worksheet {sheet.title!r} exceeds the "
                f"{_MAX_MATERIALIZED_SPREADSHEET_CELLS:,}-cell review limit"
            )
            reviewed_cells = materialized_cells[
                :_MAX_MATERIALIZED_SPREADSHEET_CELLS
            ]
        else:
            reviewed_cells = materialized_cells
        nonempty = sum(cell.value not in (None, "") for cell in reviewed_cells)
        formulas = sum(
            isinstance(cell.value, str) and cell.value.startswith("=")
            for cell in reviewed_cells
        )
        if nonempty == 0:
            warnings.append(f"worksheet {sheet.title!r} is empty")
        sheets.append({
            "name": sheet.title,
            "rows": int(sheet.max_row),
            "columns": int(sheet.max_column),
            "materialized_cells": len(materialized_cells),
            "nonempty_cells": nonempty,
            "formulas": formulas,
            "charts": len(sheet._charts),
            "tables": len(sheet.tables),
            "merged_ranges": len(sheet.merged_cells.ranges),
            "freeze_panes": str(sheet.freeze_panes) if sheet.freeze_panes else None,
            "auto_filter": str(sheet.auto_filter.ref) if sheet.auto_filter.ref else None,
        })
    workbook.close()
    if not sheets:
        errors.append("spreadsheet contains no worksheets")
    return ({"sheet_count": len(sheets), "sheets": sheets}, errors, warnings)


def _pdf(blob: bytes) -> tuple[dict[str, Any], list[str], list[str]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(blob))
    errors: list[str] = []
    warnings: list[str] = []
    if reader.is_encrypted:
        errors.append("encrypted PDF files are not supported")
        return ({"page_count": len(reader.pages), "encrypted": True}, errors, warnings)
    text_chars = []
    for page in reader.pages[:20]:
        try:
            text_chars.append(len((page.extract_text() or "").strip()))
        except Exception:
            text_chars.append(0)
    if not reader.pages:
        errors.append("PDF contains no pages")
    if reader.pages and not any(text_chars):
        warnings.append("the sampled PDF pages contain no extractable text")
    return ({
        "page_count": len(reader.pages),
        "sampled_page_text_characters": text_chars,
        "encrypted": False,
    }, errors, warnings)


def _textual(path: Path, blob: bytes) -> tuple[dict[str, Any], list[str], list[str]]:
    errors: list[str] = []
    warnings: list[str] = []
    try:
        text = blob.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise DocumentReviewError("text document must be UTF-8") from exc
    suffix = path.suffix.lower()
    details: dict[str, Any] = {
        "characters": len(text),
        "lines": len(text.splitlines()),
    }
    if not text.strip():
        errors.append("document is empty")
    if suffix in {".csv", ".tsv"}:
        delimiter = "\t" if suffix == ".tsv" else ","
        rows = list(csv.reader(io.StringIO(text), delimiter=delimiter))
        details.update({
            "rows": len(rows),
            "columns": max((len(row) for row in rows), default=0),
            "header": rows[0] if rows else [],
        })
    elif suffix in {".md", ".markdown"}:
        headings = [line.strip() for line in text.splitlines() if line.lstrip().startswith("#")]
        details["headings"] = headings[:100]
    elif suffix in {".html", ".htm"}:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(text, "html.parser")
        details.update({
            "title": soup.title.string.strip() if soup.title and soup.title.string else None,
            "headings": [node.get_text(" ", strip=True) for node in soup.select("h1,h2,h3")][:100],
            "scripts": len(soup.find_all("script")),
            "external_resources": len(soup.select('[src^="http"],[href^="http"]')),
        })
    elif suffix == ".svg":
        from lxml import etree

        if "<!DOCTYPE" in text.upper() or "<!ENTITY" in text.upper():
            errors.append("SVG must not contain DOCTYPE or ENTITY declarations")
        parser = etree.XMLParser(resolve_entities=False, no_network=True)
        root = etree.fromstring(blob, parser=parser)
        details.update({
            "root": etree.QName(root).localname,
            "view_box": root.get("viewBox"),
            "elements": sum(1 for _ in root.iter()),
        })
    return details, errors, warnings


def review_document(path: str) -> dict[str, Any]:
    """Return structural evidence and deterministic diagnostics for a file."""

    resolved = _resolve_workspace_file(path)
    blob = resolved.read_bytes()
    suffix = resolved.suffix.lower()
    if suffix == ".pptx":
        details, errors, warnings = _presentation(blob)
        format_name = "pptx"
    elif suffix == ".docx":
        details, errors, warnings = _word(blob)
        format_name = "docx"
    elif suffix == ".xlsx":
        details, errors, warnings = _spreadsheet(blob)
        format_name = "xlsx"
    elif suffix == ".pdf":
        details, errors, warnings = _pdf(blob)
        format_name = "pdf"
    elif suffix in {".md", ".markdown", ".html", ".htm", ".csv", ".tsv", ".svg", ".txt"}:
        details, errors, warnings = _textual(resolved, blob)
        format_name = suffix.removeprefix(".")
    else:
        raise DocumentReviewError(
            f"unsupported document type: {suffix or '(no extension)'}"
        )
    return {
        "path": path,
        "format": format_name,
        "size_bytes": len(blob),
        "source_hash": "sha256:" + hashlib.sha256(blob).hexdigest(),
        "valid": not errors,
        "errors": errors,
        "warnings": warnings,
        "details": details,
    }


__all__ = ["DocumentReviewError", "review_document"]
