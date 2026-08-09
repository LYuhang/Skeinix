"""PPTX parser with slide-aware citation metadata."""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from .archive import validate_office_archive
from .base import EmptyDocumentError, ParsedSegment, ParseError, Parser


class PptxParser(Parser):
    def parse(self, blob: bytes) -> list[ParsedSegment]:
        validate_office_archive(blob, required_prefix="ppt/")
        try:
            presentation = Presentation(BytesIO(blob))
        except Exception as exc:
            raise ParseError(f"Cannot parse pptx: {exc}") from exc
        segments: list[ParsedSegment] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = getattr(shape, "text", "").strip()
                    if text:
                        parts.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        line = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if line:
                            parts.append(line)
            if parts:
                segments.append(
                    ParsedSegment(text="\n".join(parts), metadata={"slide": slide_number})
                )
        if not segments:
            raise EmptyDocumentError("pptx has no readable content")
        return segments
