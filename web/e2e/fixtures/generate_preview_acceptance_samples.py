#!/usr/bin/env python3
"""Generate small, real documents for the Preview browser acceptance test."""

from __future__ import annotations

import argparse
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas


def _write_docx(target: Path) -> None:
    document = Document()
    document.add_heading("DOCX acceptance marker", level=1)
    document.add_paragraph("Rendered by the product's Word preview.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "name"
    table.cell(0, 1).text = "status"
    table.cell(1, 0).text = "Preview"
    table.cell(1, 1).text = "verified"
    document.save(target)


def _write_pptx(target: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX acceptance marker"
    slide.placeholders[1].text = "Rendered by the product's PowerPoint preview."
    presentation.save(target)


def _write_xlsx(target: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Acceptance"
    sheet.append(["name", "status"])
    sheet.append(["XLSX acceptance marker", "verified"])
    workbook.save(target)


def _write_pdf(target: Path) -> None:
    document = canvas.Canvas(str(target))
    document.drawString(72, 760, "PDF acceptance marker")
    document.showPage()
    document.save()


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("output_dir", type=Path)
    args = parser.parse_args()
    output_dir = args.output_dir.resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    docx = output_dir / "acceptance.docx"
    _write_docx(docx)
    _write_pdf(output_dir / "acceptance.pdf")
    _write_pptx(output_dir / "acceptance.pptx")
    _write_xlsx(output_dir / "acceptance.xlsx")
    (output_dir / "acceptance.csv").write_text(
        "name,status\nCSV acceptance marker,verified\n",
        encoding="utf-8",
    )
    (output_dir / "acceptance.tsv").write_text(
        "name\tstatus\nTSV acceptance marker\tverified\n",
        encoding="utf-8",
    )
    (output_dir / "acceptance.jsonl").write_text(
        '{"name":"JSONL acceptance marker","status":"verified"}\n',
        encoding="utf-8",
    )
    (output_dir / "acceptance.txt").write_bytes(
        b"\xef\xbb\xbfText acceptance marker\r\nOriginal line\r\n"
    )
    (output_dir / "acceptance.rtf").write_bytes(
        b"{\\rtf1\\ansi\\deff0{\\fonttbl{\\f0 Liberation Sans;}}"
        b"\\f0\\fs28 Legacy Office acceptance marker\\par}"
    )


if __name__ == "__main__":
    main()
